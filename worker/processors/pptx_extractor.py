def extract(file_bytes: bytes) -> str:
    from pptx import Presentation
    import io

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_lines = [f"Slide {i}:"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    slide_lines.append(shape.text.strip())
            if len(slide_lines) > 1:
                slides.append("\n".join(slide_lines))
        result = "\n\n".join(slides).strip()
        if not result:
            raise ValueError("No readable text found in PPTX")
        return result
    except Exception as exc:
        raise Exception(f"PPTX extraction failed: {exc}") from exc
